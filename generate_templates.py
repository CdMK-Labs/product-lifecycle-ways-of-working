"""
Generates all Office template files for the product lifecycle stages.
Run once: python3 generate_templates.py
"""

import os
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
from openpyxl.utils import get_column_letter
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────
DARK_GREEN  = "02462F"
MID_GREEN   = "0D6A4B"
TEAL        = "0B885A"
LIGHT_GREEN = "82CE71"
GRAY        = "999999"
WHITE       = "FFFFFF"
NEAR_BLACK  = "1A1A1A"
LIGHT_BG    = "F5FAF5"

# ── Helpers ─────────────────────────────────────────────────────────────────

def ensure(path):
    os.makedirs(path, exist_ok=True)

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def ppt_rgb(hex_str):
    h = hex_str.lstrip("#")
    return PptRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def xl_fill(hex_str):
    return PatternFill("solid", fgColor=hex_str)

def xl_font(bold=False, size=11, color=NEAR_BLACK, name="Calibri"):
    return Font(bold=bold, size=size, color=color, name=name)

def xl_border():
    s = Side(border_style="thin", color="CCCCCC")
    return Border(left=s, right=s, top=s, bottom=s)

def xl_center():
    return Alignment(horizontal="center", vertical="center", wrap_text=True)

def xl_left():
    return Alignment(horizontal="left", vertical="center", wrap_text=True)

def xl_header(ws, row, cols, bg=DARK_GREEN):
    for col_idx, (label, width) in enumerate(cols, 1):
        c = ws.cell(row=row, column=col_idx, value=label)
        c.font = xl_font(bold=True, color=WHITE)
        c.fill = xl_fill(bg)
        c.alignment = xl_center()
        c.border = xl_border()
        ws.column_dimensions[get_column_letter(col_idx)].width = width

def xl_row(ws, row, values, shade=False):
    bg = "EEF6EE" if shade else WHITE
    for col_idx, val in enumerate(values, 1):
        c = ws.cell(row=row, column=col_idx, value=val)
        c.fill = xl_fill(bg)
        c.alignment = xl_left()
        c.border = xl_border()
        c.font = xl_font()

def xl_title(ws, title, subtitle=""):
    ws.merge_cells("A1:H1")
    t = ws["A1"]
    t.value = title
    t.font = xl_font(bold=True, size=16, color=WHITE)
    t.fill = xl_fill(DARK_GREEN)
    t.alignment = xl_center()
    ws.row_dimensions[1].height = 32
    if subtitle:
        ws.merge_cells("A2:H2")
        s = ws["A2"]
        s.value = subtitle
        s.font = xl_font(size=10, color=GRAY)
        s.fill = xl_fill("F0F0F0")
        s.alignment = xl_center()
        ws.row_dimensions[2].height = 18

def xl_section(ws, row, label, span=8):
    ws.merge_cells(f"A{row}:{get_column_letter(span)}{row}")
    c = ws[f"A{row}"]
    c.value = label
    c.font = xl_font(bold=True, size=11, color=WHITE)
    c.fill = xl_fill(MID_GREEN)
    c.alignment = xl_left()
    ws.row_dimensions[row].height = 22

# ── Word helpers ─────────────────────────────────────────────────────────────

def doc_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    run = p.runs[0] if p.runs else p.add_run(text)
    run.font.color.rgb = rgb(DARK_GREEN if level == 1 else MID_GREEN)
    return p

def doc_para(doc, text, size=11):
    p = doc.add_paragraph(text)
    for run in p.runs:
        run.font.size = Pt(size)
    return p

def doc_field(doc, label, placeholder=""):
    p = doc.add_paragraph()
    run_label = p.add_run(f"{label}: ")
    run_label.bold = True
    run_label.font.color.rgb = rgb(MID_GREEN)
    run_val = p.add_run(placeholder or f"[{label}]")
    run_val.font.color.rgb = rgb(GRAY)
    return p

def doc_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    # Header row
    hrow = table.rows[0]
    for i, h in enumerate(headers):
        cell = hrow.cells[i]
        cell.text = h
        run = cell.paragraphs[0].runs[0]
        run.bold = True
        run.font.color.rgb = rgb(WHITE)
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), DARK_GREEN)
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:val"), "clear")
        tcPr.append(shd)
    # Data rows
    for r_idx, row_data in enumerate(rows):
        row = table.rows[r_idx + 1]
        for c_idx, val in enumerate(row_data):
            row.cells[c_idx].text = str(val)
    return table

def doc_checklist(doc, items):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run("☐  ").font.color.rgb = rgb(MID_GREEN)
        p.add_run(item)

def doc_meta(doc, product="[Product / Feature]", date="[DD/MM/YYYY]", author="[Name, Role]"):
    doc_field(doc, "Product / Feature", product)
    doc_field(doc, "Date", date)
    doc_field(doc, "Author / Owner", author)
    doc.add_paragraph()

# ── PowerPoint helpers ────────────────────────────────────────────────────────

def ppt_blank(prs):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)

def ppt_text_box(slide, text, l, t, w, h, size=24, bold=False,
                 color=NEAR_BLACK, align=PP_ALIGN.LEFT, bg=None):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)
    if bg:
        from pptx.oxml.ns import qn as pqn
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = ppt_rgb(bg)
    return txBox

def ppt_rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(fill_color)
    if line_color:
        shape.line.color.rgb = ppt_rgb(line_color)
    else:
        shape.line.fill.background()
    return shape

def ppt_accent_bar(slide, color=DARK_GREEN, width=10):
    ppt_rect(slide, 0, 0, width, 0.12, color)

def ppt_title_slide(prs, title, subtitle="Product Lifecycle"):
    slide = ppt_blank(prs)
    ppt_rect(slide, 0, 0, 10, 7.5, LIGHT_BG)
    ppt_rect(slide, 0, 0, 10, 0.15, DARK_GREEN)
    ppt_rect(slide, 0, 7.35, 10, 0.15, DARK_GREEN)
    ppt_text_box(slide, subtitle, 0.5, 1.8, 9, 0.5, size=14,
                 color=MID_GREEN, align=PP_ALIGN.CENTER)
    ppt_text_box(slide, title, 0.5, 2.4, 9, 1.5, size=36, bold=True,
                 color=DARK_GREEN, align=PP_ALIGN.CENTER)
    ppt_text_box(slide, "[Product / Feature]    |    [Date]    |    [Author]",
                 0.5, 6.6, 9, 0.5, size=11, color=GRAY, align=PP_ALIGN.CENTER)

def ppt_content_slide(prs, title, bullets, color=MID_GREEN):
    slide = ppt_blank(prs)
    ppt_accent_bar(slide, DARK_GREEN)
    ppt_text_box(slide, title, 0.4, 0.25, 9.2, 0.6, size=20, bold=True, color=DARK_GREEN)
    # horizontal rule
    ppt_rect(slide, 0.4, 0.92, 9.2, 0.02, MID_GREEN)
    y = 1.05
    for b in bullets:
        ppt_text_box(slide, f"–  {b}", 0.5, y, 8.8, 0.45, size=13, color=NEAR_BLACK)
        y += 0.42

def ppt_table_slide(prs, title, headers, rows):
    slide = ppt_blank(prs)
    ppt_accent_bar(slide, DARK_GREEN)
    ppt_text_box(slide, title, 0.4, 0.25, 9.2, 0.6, size=20, bold=True, color=DARK_GREEN)
    ppt_rect(slide, 0.4, 0.92, 9.2, 0.02, MID_GREEN)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table = slide.shapes.add_table(n_rows, n_cols,
                                   Inches(0.4), Inches(1.05),
                                   Inches(9.2), Inches(min(0.4 * n_rows, 5.8))).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ppt_rgb(DARK_GREEN)
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = ppt_rgb(WHITE)
        p.runs[0].font.size = PptPt(11)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ppt_rgb("EEF6EE")
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].font.size = PptPt(10)

# ════════════════════════════════════════════════════════════════════════════
#  FILE GENERATORS
# ════════════════════════════════════════════════════════════════════════════

# ── 1. INFLOW ────────────────────────────────────────────────────────────────

def inflow_request(path):
    doc = Document()
    doc_heading(doc, "Inflow Request")
    doc_meta(doc)
    doc_heading(doc, "Problem Statement", 2)
    doc_para(doc, "Describe the business need, opportunity, or problem in 2–3 sentences.")
    doc_para(doc, "[Your text here]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Proposed Solution or Direction", 2)
    doc_para(doc, "Optional. Leave blank if unknown.")
    doc_para(doc, "[Your text here]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Strategic Alignment", 2)
    doc_field(doc, "OKR or strategic objective", "")
    doc_heading(doc, "Expected Benefit", 2)
    doc_para(doc, "[Your text here]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Urgency", 2)
    doc_checklist(doc, ["Critical — needs attention immediately",
                         "High — should be picked up in the next cycle",
                         "Medium — important but not time-sensitive",
                         "Low — nice to have"])
    doc_heading(doc, "Proposed Owner", 2)
    doc_field(doc, "Product Owner", "")
    doc_field(doc, "Key Stakeholders", "")
    doc.save(path)

def stakeholder_brief(path):
    doc = Document()
    doc_heading(doc, "Stakeholder Brief")
    doc_meta(doc)
    doc_heading(doc, "Overview", 2)
    doc_para(doc, "[One paragraph describing what this initiative is about and why it matters.]"
             ).runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Stakeholder Map", 2)
    doc_table(doc,
              ["Stakeholder", "Role", "Interest / Concern", "Engagement Level"],
              [["[Name]", "[Role]", "[What they care about]", "Inform / Consult / Involve"],
               ["[Name]", "[Role]", "[What they care about]", "Inform / Consult / Involve"]])
    doc.add_paragraph()
    doc_heading(doc, "Communication Plan", 2)
    doc_table(doc,
              ["Audience", "Message", "Channel", "Frequency"],
              [["[Group]", "[Key message]", "[Email / Slack / Meeting]", "[Weekly / Monthly]"]])
    doc.add_paragraph()
    doc_heading(doc, "Risks and Concerns", 2)
    doc_para(doc, "[Summarise any stakeholder concerns that need to be managed.]"
             ).runs[0].font.color.rgb = rgb(GRAY)
    doc.save(path)

# ── 2. VALUE PROPOSITION ─────────────────────────────────────────────────────

def vp_canvas(path):
    doc = Document()
    doc_heading(doc, "Value Proposition Canvas")
    doc_meta(doc)
    doc_heading(doc, "Customer Profile", 2)
    doc_heading(doc, "Customer Jobs", 3)
    doc_para(doc, "What tasks, problems, or needs is the customer trying to address?")
    for i in range(1, 4):
        doc_para(doc, f"[Job {i}]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Customer Pains", 3)
    for i in range(1, 4):
        doc_para(doc, f"[Pain {i}]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Customer Gains", 3)
    for i in range(1, 4):
        doc_para(doc, f"[Gain {i}]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Value Map", 2)
    doc_heading(doc, "Products & Services", 3)
    for i in range(1, 3):
        doc_para(doc, f"[Product / Feature {i}]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Pain Relievers", 3)
    for i in range(1, 3):
        doc_para(doc, f"[Pain Reliever {i}] → addresses [Pain X]"
                 ).runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Gain Creators", 3)
    for i in range(1, 3):
        doc_para(doc, f"[Gain Creator {i}] → delivers [Gain X]"
                 ).runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Fit Assessment", 2)
    doc_checklist(doc, ["Strong fit — core pains and gains are addressed",
                         "Partial fit — some gaps remain; further discovery needed",
                         "Weak fit — significant rethink required"])
    doc.save(path)

def success_metrics(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Success Metrics"
    xl_title(ws, "Success Metrics Template",
             "Define and agree measurable outcomes before moving to Product Discovery")
    ws.row_dimensions[3].height = 6
    # Meta
    for r, (label, val) in enumerate([("Product / Feature", "[Name]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Owner", "[Name, Role]"),
                                       ("Linked OKR", "[OKR reference]")], start=4):
        ws.cell(r, 1).value = label
        ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val
        ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[8].height = 10
    # Primary metric
    xl_section(ws, 9, "Primary Metric — The single most important indicator of success")
    xl_header(ws, 10, [("Metric", 30), ("Definition", 35), ("Baseline", 16),
                        ("Target", 16), ("Data Source", 25), ("Owner", 20)])
    xl_row(ws, 11, ["[Metric name]", "[How it is calculated]",
                     "[Current]", "[Target]", "[Source]", "[Name]"])
    ws.row_dimensions[12].height = 8
    # Supporting metrics
    xl_section(ws, 13, "Supporting Metrics")
    xl_header(ws, 14, [("Metric", 30), ("Definition", 35), ("Baseline", 16),
                        ("Target", 16), ("Data Source", 25), ("Owner", 20)])
    for r in range(15, 18):
        xl_row(ws, r, ["[Metric name]", "[Definition]",
                        "[Current]", "[Target]", "[Source]", "[Name]"], r % 2 == 0)
    ws.row_dimensions[18].height = 8
    # Counter metrics
    xl_section(ws, 19, "Counter Metrics — Monitor for unintended negative effects")
    xl_header(ws, 20, [("Metric", 30), ("Definition", 35), ("Acceptable Threshold", 26),
                        ("Data Source", 25), ("Owner", 20)])
    xl_row(ws, 21, ["[Metric name]", "[Definition]",
                     "[Threshold]", "[Source]", "[Name]"])
    ws.row_dimensions[22].height = 8
    # Reporting
    xl_section(ws, 23, "Reporting")
    xl_header(ws, 24, [("Report", 30), ("Audience", 25), ("Format", 20),
                        ("Cadence", 20)])
    xl_row(ws, 25, ["[Report name]", "[Audience]", "[Dashboard / Doc]", "[Weekly / Monthly]"])
    ws.sheet_view.showGridLines = False
    wb.save(path)

# ── 3. PRODUCT DISCOVERY ─────────────────────────────────────────────────────

def research_plan(path):
    doc = Document()
    doc_heading(doc, "Research Plan")
    doc_meta(doc)
    doc_field(doc, "Lead Researcher", "")
    doc.add_paragraph()
    doc_heading(doc, "Research Questions", 2)
    doc_para(doc, "What do we need to learn? Frame as questions, not assumptions.")
    for i in range(1, 4):
        doc_para(doc, f"{i}.  [Research question {i}]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Hypotheses", 2)
    doc_table(doc,
              ["Hypothesis", "Confidence", "How We Will Test It"],
              [["We believe that [X]", "Low / Med / High", "[Method]"],
               ["We believe that [Y]", "Low / Med / High", "[Method]"]])
    doc.add_paragraph()
    doc_heading(doc, "Research Methods", 2)
    doc_table(doc,
              ["Method", "Purpose", "Participants", "Timeline"],
              [["User interviews", "[Why]", "N = [X], [profile]", "[Dates]"],
               ["Survey", "[Why]", "N = [X], [profile]", "[Dates]"],
               ["Usability testing", "[Why]", "N = [X], [profile]", "[Dates]"]])
    doc.add_paragraph()
    doc_heading(doc, "Outputs", 2)
    doc_table(doc,
              ["Deliverable", "Owner", "Due Date"],
              [["Research synthesis document", "[Name]", "[Date]"],
               ["Key insights presentation", "[Name]", "[Date]"],
               ["Recommendation and next steps", "[Name]", "[Date]"]])
    doc.save(path)

def assumption_log(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Assumption Log"
    xl_title(ws, "Assumption Log",
             "Capture, rate, and prioritise assumptions for testing")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product / Feature", "[Name]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Owner", "[Name, Role]")], start=4):
        ws.cell(r, 1).value = label
        ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val
        ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[7].height = 10
    xl_section(ws, 8, "Assumptions  —  Priority: High Importance × Low Confidence", span=8)
    xl_header(ws, 9, [("ID", 6), ("Assumption", 45), ("Category", 20),
                       ("Importance", 14), ("Confidence", 14),
                       ("Status", 18), ("Evidence / Notes", 35)])
    statuses = ["Untested", "Validated", "Invalidated"]
    categories = ["User", "Market", "Technical", "Business"]
    for i in range(1, 8):
        shade = i % 2 == 0
        xl_row(ws, 9 + i,
               [f"A{i:02d}", f"[Assumption {i}]",
                categories[(i - 1) % 4], "High / Med / Low",
                "High / Med / Low", statuses[0], ""], shade)
    ws.row_dimensions[17].height = 10
    xl_section(ws, 18, "Priority Testing Queue — High importance, Low confidence", span=8)
    xl_header(ws, 19, [("ID", 6), ("Assumption Summary", 45), ("Test Method", 30),
                        ("Owner", 20), ("Target Date", 16), ("Result", 30)])
    for i in range(20, 23):
        xl_row(ws, i, ["[A##]", "[Assumption]", "[Method]", "[Name]", "[Date]", ""], i % 2 == 0)
    ws.sheet_view.showGridLines = False
    wb.save(path)

# ── 4. ALIGN & PLAN ──────────────────────────────────────────────────────────

def delivery_plan(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Delivery Plan"
    xl_title(ws, "Delivery Plan",
             "Scope, sequence, capacity, and milestones for the delivery increment")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product / Feature", "[Name]"),
                                       ("Product Owner", "[Name]"),
                                       ("Tech Lead", "[Name]"),
                                       ("Date", "[DD/MM/YYYY]")], start=4):
        ws.cell(r, 1).value = label; ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val;   ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[8].height = 10
    # Scope
    xl_section(ws, 9, "Scope", span=6)
    xl_header(ws, 10, [("In Scope", 50), ("Out of Scope", 50)])
    ws.merge_cells("A10:A10"); ws.merge_cells("B10:B10")
    for r in range(11, 15):
        xl_row(ws, r, ["[In scope item]", "[Out of scope item]"], r % 2 == 0)
    ws.row_dimensions[15].height = 10
    # Sequence
    xl_section(ws, 16, "Delivery Sequence & Estimation", span=6)
    xl_header(ws, 17, [("Phase", 18), ("Description", 38), ("Team", 18),
                        ("Start", 14), ("End", 14), ("Estimate (days)", 18), ("Dependencies", 28)])
    for i in range(1, 5):
        xl_row(ws, 17 + i, [f"Phase {i}", "[Description]", "[Team]",
                              "[Date]", "[Date]", "", "[None / Phase X]"], i % 2 == 0)
    ws.row_dimensions[22].height = 10
    # Team
    xl_section(ws, 23, "Team & Capacity", span=6)
    xl_header(ws, 24, [("Role", 22), ("Name", 22), ("Availability (days/week)", 28),
                        ("Notes", 40)])
    for role in ["Product Owner", "Tech Lead", "Engineer", "Designer"]:
        xl_row(ws, 25 + ["Product Owner", "Tech Lead", "Engineer", "Designer"].index(role),
               [role, "[Name]", "[X]", ""])
    ws.row_dimensions[29].height = 10
    # Milestones
    xl_section(ws, 30, "Key Milestones", span=6)
    xl_header(ws, 31, [("Milestone", 40), ("Date", 16), ("Owner", 22), ("Status", 18)])
    for i in range(1, 5):
        xl_row(ws, 31 + i, [f"[Milestone {i}]", "[Date]", "[Name]", "Planned"], i % 2 == 0)
    ws.sheet_view.showGridLines = False
    wb.save(path)

def risk_register(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Risk Register"
    xl_title(ws, "Risk Register",
             "Score = Likelihood (1–5) × Impact (1–5) | 1–10 Low · 11–15 Med · 16–25 High")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product / Feature", "[Name]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Owner", "[Name, Role]")], start=4):
        ws.cell(r, 1).value = label; ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val;   ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[7].height = 10
    xl_section(ws, 8, "Risk Log", span=8)
    xl_header(ws, 9, [("ID", 7), ("Risk Description", 40), ("Category", 18),
                       ("Likelihood\n1–5", 14), ("Impact\n1–5", 14), ("Score", 10),
                       ("Owner", 18), ("Mitigation / Response", 38), ("Status", 16)])
    categories = ["Technical", "People", "Process", "External"]
    for i in range(1, 8):
        shade = i % 2 == 0
        xl_row(ws, 9 + i,
               [f"R{i:02d}", f"[Risk description {i}]",
                categories[(i - 1) % 4], "", "", "", "[Name]",
                "[Mitigation action]", "Open"], shade)
    ws.row_dimensions[17].height = 10
    xl_section(ws, 18, "Escalated Risks — Score ≥ 16", span=8)
    xl_header(ws, 19, [("ID", 7), ("Risk", 40), ("Score", 10),
                        ("Escalated To", 22), ("Date", 14), ("Resolution", 38)])
    xl_row(ws, 20, ["[R##]", "[Description]", "", "[Name, Role]", "[Date]", ""])
    ws.sheet_view.showGridLines = False
    wb.save(path)

# ── 5. DEVELOP & DELIVER ─────────────────────────────────────────────────────

def definition_of_done(path):
    doc = Document()
    doc_heading(doc, "Definition of Done")
    doc_field(doc, "Product / Team", "")
    doc_field(doc, "Date Agreed", "")
    doc_field(doc, "Review Date", "")
    doc.add_paragraph()
    doc_para(doc,
             "The Definition of Done is a shared agreement on the quality standards "
             "a piece of work must meet before it can be considered complete.")
    sections = {
        "Code Quality": [
            "Code reviewed and approved by at least one peer",
            "No critical or high-severity linting errors",
            "Unit tests written and passing (coverage ≥ [X]%)",
            "Integration tests passing",
            "No known bugs introduced by the change",
        ],
        "Functionality": [
            "Feature meets acceptance criteria defined in the user story",
            "Edge cases and error states are handled",
            "Tested on [browser / device / environment list]",
            "Accessibility requirements met ([WCAG level])",
        ],
        "Documentation": [
            "Code comments added where logic is non-obvious",
            "API documentation updated (if applicable)",
            "Runbooks or operational guides updated (if applicable)",
            "Release notes entry drafted",
        ],
        "Deployment Readiness": [
            "Feature flag configured (if applicable)",
            "Environment variables and secrets managed",
            "Monitoring and alerting updated",
            "Performance benchmarks checked",
        ],
        "Sign-off": [
            "Product Owner has reviewed and accepted the feature",
            "QA sign-off obtained (if applicable)",
        ],
    }
    for section, items in sections.items():
        doc_heading(doc, section, 2)
        doc_checklist(doc, items)
    doc.save(path)

def sprint_review_pptx(path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    ppt_title_slide(prs, "Sprint Review", "Product Lifecycle — Develop & Deliver")
    ppt_content_slide(prs, "Sprint Goal & Summary",
                      ["Sprint: [Number / Name]    |    Dates: [Start] – [End]",
                       "Goal: [What was the team trying to achieve this sprint?]",
                       "Overall status:  ✅ On track  /  ⚠️ At risk  /  ❌ Off track"])
    ppt_table_slide(prs, "Completed Work",
                    ["Item / Story", "Description", "Demo?", "Notes"],
                    [["[Story ID]", "[Brief description]", "Yes", ""],
                     ["[Story ID]", "[Brief description]", "Yes", ""],
                     ["[Story ID]", "[Brief description]", "No — carried forward", ""]])
    ppt_content_slide(prs, "Stakeholder Feedback",
                      ["[Feedback point 1 from demo]",
                       "[Feedback point 2 from demo]",
                       "[Feedback point 3 from demo]",
                       "Actions raised: [List or 'None']"])
    ppt_table_slide(prs, "Metrics Update",
                    ["Metric", "Previous", "Current", "Trend"],
                    [["[Metric name]", "[Value]", "[Value]", "↑ / → / ↓"],
                     ["[Metric name]", "[Value]", "[Value]", "↑ / → / ↓"]])
    ppt_content_slide(prs, "Next Steps",
                      ["Next sprint planning: [Date]",
                       "Key focus: [Topic]",
                       "Dependencies to resolve: [List]",
                       "Risks to watch: [List]"])
    prs.save(path)

# ── 6. DEPLOY ────────────────────────────────────────────────────────────────

def deployment_checklist(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Deployment Checklist"
    xl_title(ws, "Deployment Checklist",
             "Complete all items before and after each production deployment")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product / Release", "[Name and version]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Deployment Lead", "[Name]"),
                                       ("Approver", "[Name, Role]")], start=4):
        ws.cell(r, 1).value = label; ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val;   ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[8].height = 10
    def checklist_section(start_row, title, items):
        xl_section(ws, start_row, title, span=4)
        xl_header(ws, start_row + 1,
                  [("Done", 8), ("Item", 60), ("Owner", 20), ("Notes", 28)])
        for i, item in enumerate(items):
            xl_row(ws, start_row + 2 + i,
                   ["☐", item, "[Name]", ""], (start_row + 2 + i) % 2 == 0)
        return start_row + 2 + len(items) + 1

    row = checklist_section(9, "Pre-Deployment — Code & Quality Gates", [
        "All items in scope have met the Definition of Done",
        "All automated tests passing in staging",
        "Security scan completed with no critical findings",
        "Performance benchmarks within acceptable thresholds",
        "Dependency versions pinned and reviewed",
    ])
    row = checklist_section(row, "Pre-Deployment — Environment Readiness", [
        "Staging deployment tested and verified",
        "Production environment capacity confirmed",
        "Feature flags configured correctly",
        "Environment variables and secrets updated",
        "Database migrations prepared and tested",
    ])
    row = checklist_section(row, "Deployment", [
        "Deployment initiated",
        "Deployment completed without errors",
        "Database migrations executed successfully",
        "Health checks passing",
        "Smoke tests passed",
    ])
    row = checklist_section(row, "Post-Deployment", [
        "Key user journeys verified in production",
        "Monitoring dashboards checked — no anomalies",
        "Error rates within normal range",
        "Performance metrics within acceptable range",
        "Stakeholders notified of successful deployment",
        "Release notes published",
    ])
    ws.sheet_view.showGridLines = False
    wb.save(path)

def rollback_plan(path):
    doc = Document()
    doc_heading(doc, "Rollback Plan")
    doc_field(doc, "Product / Release", "")
    doc_field(doc, "Date", "")
    doc_field(doc, "Owner", "")
    doc.add_paragraph()
    doc_heading(doc, "Decision Authority", 2)
    doc_table(doc,
              ["Role", "Name", "Contact"],
              [["Primary", "[Name]", "[Slack / Phone]"],
               ["Backup", "[Name]", "[Slack / Phone]"]])
    doc.add_paragraph()
    doc_heading(doc, "Rollback Triggers", 2)
    doc_para(doc, "Initiate this procedure immediately if any of the following are observed:")
    doc_checklist(doc, [
        "[Condition 1 — e.g. error rate > X% for more than Y minutes]",
        "[Condition 2 — e.g. P1 incident raised by support team]",
        "[Condition 3 — e.g. key metric drops below threshold]",
    ])
    doc_heading(doc, "Rollback Steps", 2)
    doc_table(doc,
              ["Step", "Action", "Owner", "Est. Time"],
              [["1", "Confirm rollback decision", "[Name]", "2 min"],
               ["2", "Notify on-call team and stakeholders", "[Name]", "5 min"],
               ["3", "[Technical rollback — e.g. revert via CI/CD]", "[Name]", "[X min]"],
               ["4", "Database rollback if required", "[Name]", "[X min]"],
               ["5", "Verify health checks pass", "[Name]", "[X min]"],
               ["6", "Confirm error rates normalising", "[Name]", "[X min]"],
               ["7", "Notify stakeholders of completion", "[Name]", "5 min"]])
    doc.add_paragraph()
    doc_heading(doc, "Post-Rollback Actions", 2)
    doc_checklist(doc, [
        "Incident report created",
        "Root cause identified",
        "Fix validated in staging before re-attempting deployment",
        "Retrospective scheduled",
    ])
    doc.save(path)

# ── 7. OPERATE & MONITOR ─────────────────────────────────────────────────────

def incident_report(path):
    doc = Document()
    doc_heading(doc, "Incident Report")
    doc_field(doc, "Incident ID", "[INC-XXXX]")
    doc_field(doc, "Product / Service", "")
    doc_field(doc, "Severity", "P1 / P2 / P3 / P4")
    doc_field(doc, "Date", "")
    doc_field(doc, "Report Author", "")
    doc.add_paragraph()
    doc_heading(doc, "Summary", 2)
    doc_para(doc, "[One paragraph: what happened, the impact, and how it was resolved.]"
             ).runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Timeline", 2)
    doc_table(doc,
              ["Time (UTC)", "Event"],
              [["[HH:MM]", "Incident first detected / alert fired"],
               ["[HH:MM]", "On-call engineer notified"],
               ["[HH:MM]", "Investigation begun"],
               ["[HH:MM]", "Root cause identified"],
               ["[HH:MM]", "Mitigation applied"],
               ["[HH:MM]", "Service restored"],
               ["[HH:MM]", "Incident closed"]])
    doc.add_paragraph()
    doc_heading(doc, "Impact", 2)
    doc_field(doc, "Users affected", "")
    doc_field(doc, "Services affected", "")
    doc_field(doc, "Data impact", "None")
    doc_field(doc, "SLA breach", "Yes / No")
    doc_field(doc, "Total duration", "")
    doc_heading(doc, "Root Cause", 2)
    doc_para(doc, "[What was the underlying technical or process cause?]"
             ).runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Resolution", 2)
    doc_para(doc, "[What was done to restore service?]").runs[0].font.color.rgb = rgb(GRAY)
    doc_heading(doc, "Action Items", 2)
    doc_table(doc,
              ["Action", "Owner", "Due Date", "Status"],
              [["[Preventive action]", "[Name]", "[Date]", "Open"],
               ["[Process improvement]", "[Name]", "[Date]", "Open"]])
    doc.save(path)

def operational_review_pptx(path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    ppt_title_slide(prs, "Operational Review", "Product Lifecycle — Operate & Monitor")
    ppt_table_slide(prs, "Service Health",
                    ["Metric", "Target", "Actual", "Status"],
                    [["Availability", "[X]%", "[X]%", "✅ / ⚠️ / ❌"],
                     ["P95 Response Time", "[Xms]", "[Xms]", "✅ / ⚠️ / ❌"],
                     ["Error Rate", "< [X]%", "[X]%", "✅ / ⚠️ / ❌"],
                     ["[Custom SLA]", "[Target]", "[Actual]", "✅ / ⚠️ / ❌"]])
    ppt_table_slide(prs, "Incidents This Period",
                    ["ID", "Severity", "Date", "Duration", "Root Cause", "Resolved?"],
                    [["[INC-XXX]", "P[1–4]", "[Date]", "[Xh Ym]", "[Brief cause]", "Yes / No"]])
    ppt_content_slide(prs, "Highlights & Concerns",
                      ["✅  [What went well — achievement or improvement]",
                       "✅  [What went well]",
                       "⚠️  [Issue or concern that needs attention]",
                       "⚠️  [Recurring problem or risk area]"])
    ppt_table_slide(prs, "Improvement Actions",
                    ["Action", "Owner", "Due Date", "Status"],
                    [["[Action]", "[Name]", "[Date]", "In progress / Planned"],
                     ["[Action]", "[Name]", "[Date]", "In progress / Planned"]])
    ppt_content_slide(prs, "Next Review",
                      ["Next review date: [DD/MM/YYYY]",
                       "Focus areas: [Topics to prioritise]",
                       "Attendees: [Names / Roles]"])
    prs.save(path)

# ── 8. MEASURE & LEARN ───────────────────────────────────────────────────────

def measurement_plan(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Measurement Plan"
    xl_title(ws, "Measurement Plan",
             "Define metrics, baselines, and reporting — agree before deployment")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product / Feature", "[Name]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Owner", "[Name, Role]"),
                                       ("Linked OKR / Value Proposition", "[Reference]")], start=4):
        ws.cell(r, 1).value = label; ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val;   ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[8].height = 10
    xl_section(ws, 9, "North Star Metric — Single most important indicator of user value", span=6)
    xl_header(ws, 10, [("Metric", 30), ("Definition", 40), ("Baseline", 16),
                        ("Target", 16), ("Timeframe", 18), ("Data Source", 28)])
    xl_row(ws, 11, ["[Metric name]", "[How it is calculated]",
                     "[Current]", "[Target]", "[By when]", "[Source]"])
    ws.row_dimensions[12].height = 8
    xl_section(ws, 13, "Supporting Metrics", span=6)
    xl_header(ws, 14, [("Metric", 30), ("Definition", 40), ("Baseline", 16),
                        ("Target", 16), ("Data Source", 28), ("Reporting Frequency", 22)])
    for i in range(15, 19):
        xl_row(ws, i, ["[Metric]", "[Definition]", "[Current]",
                        "[Target]", "[Source]", "Weekly / Monthly"], i % 2 == 0)
    ws.row_dimensions[19].height = 8
    xl_section(ws, 20, "Qualitative Signals", span=6)
    xl_header(ws, 21, [("Signal", 30), ("Method", 30), ("Frequency", 18), ("Owner", 22)])
    xl_row(ws, 22, ["User satisfaction", "NPS survey", "Monthly", "[Name]"])
    xl_row(ws, 23, ["Support themes", "Ticket analysis", "Weekly", "[Name]"], True)
    ws.row_dimensions[24].height = 8
    xl_section(ws, 25, "Reporting Plan", span=6)
    xl_header(ws, 26, [("Report", 30), ("Audience", 25), ("Format", 22),
                        ("Cadence", 18), ("Owner", 22)])
    xl_row(ws, 27, ["[Report name]", "[Audience]", "[Dashboard / Doc / Meeting]",
                     "[Frequency]", "[Name]"])
    ws.sheet_view.showGridLines = False
    wb.save(path)

def insights_report_pptx(path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    ppt_title_slide(prs, "Insights Report", "Product Lifecycle — Measure & Learn")
    ppt_content_slide(prs, "Executive Summary",
                      ["Period: [DD/MM/YYYY – DD/MM/YYYY]",
                       "[Key finding 1 — most important insight]",
                       "[Key finding 2]",
                       "Recommended action: [Top recommendation]"])
    ppt_table_slide(prs, "Performance Against Metrics",
                    ["Metric", "Target", "Actual", "vs. Baseline", "Trend"],
                    [["[North Star]", "[Target]", "[Actual]", "[+/- X%]", "↑ / → / ↓"],
                     ["[Supporting]", "[Target]", "[Actual]", "[+/- X%]", "↑ / → / ↓"],
                     ["[Counter]",    "[Target]", "[Actual]", "[+/- X%]", "↑ / → / ↓"]])
    ppt_table_slide(prs, "User Feedback Themes",
                    ["Theme", "Frequency", "Sentiment", "Example Quote"],
                    [["[Theme 1]", "High / Med / Low", "Positive / Neutral / Negative", '"[Quote]"'],
                     ["[Theme 2]", "High / Med / Low", "Positive / Neutral / Negative", '"[Quote]"']])
    ppt_content_slide(prs, "Key Insights",
                      ["1.  [Insight title]: [Explanation and evidence]",
                       "2.  [Insight title]: [Explanation and evidence]",
                       "3.  [Insight title]: [Explanation and evidence]"])
    ppt_table_slide(prs, "Recommendations",
                    ["Recommendation", "Priority", "Rationale", "Proposed Owner"],
                    [["[Action]", "High / Med / Low", "[Why]", "[Team]"],
                     ["[Action]", "High / Med / Low", "[Why]", "[Team]"]])
    prs.save(path)

# ── 9. ADAPT & OPTIMIZE ──────────────────────────────────────────────────────

def optimization_backlog(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Optimization Backlog"
    xl_title(ws, "Optimization Backlog",
             "Evidence-based improvement opportunities — prioritise by Impact × Confidence ÷ Effort")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product / Feature", "[Name]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Owner", "[Name, Role]")], start=4):
        ws.cell(r, 1).value = label; ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val;   ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[7].height = 10
    xl_section(ws, 8, "Backlog", span=8)
    xl_header(ws, 9, [("ID", 7), ("Opportunity", 42), ("Source", 28),
                       ("Expected Impact", 18), ("Effort", 14),
                       ("Priority", 12), ("Owner", 18), ("Status", 18)])
    for i in range(1, 9):
        shade = i % 2 == 0
        xl_row(ws, 9 + i,
               [f"O{i:02d}", f"[Opportunity {i}]",
                "[Insight / Metric / Feedback]",
                "High / Med / Low", "High / Med / Low",
                str(i), "[Name]", "Backlog"], shade)
    ws.row_dimensions[18].height = 10
    xl_section(ws, 19, "Completed Items", span=8)
    xl_header(ws, 20, [("ID", 7), ("Description", 42), ("Completed", 16),
                        ("Outcome / Result", 55)])
    for i in range(21, 24):
        xl_row(ws, i, ["[O##]", "[Description]", "[Date]",
                        "[What changed as a result]"], i % 2 == 0)
    ws.sheet_view.showGridLines = False
    wb.save(path)

def retrospective_pptx(path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    ppt_title_slide(prs, "Retrospective", "Product Lifecycle — Adapt & Optimize")
    ppt_content_slide(prs, "Start / Stop / Continue",
                      ["Format: Start — Stop — Continue",
                       "Team: [Names]    |    Date: [DD/MM/YYYY]",
                       "Facilitator: [Name]",
                       "Duration: [X minutes]"])
    ppt_table_slide(prs, "Continue — What's Working",
                    ["Item", "Votes", "Action"],
                    [["[Practice or behaviour]", "[N]", "Keep as-is / Formalise"],
                     ["[Practice or behaviour]", "[N]", "Keep as-is / Document"]])
    ppt_table_slide(prs, "Stop — What Needs to Change",
                    ["Item", "Votes", "Action"],
                    [["[Pain point or waste]", "[N]", "[How we will stop doing this]"],
                     ["[Pain point or waste]", "[N]", "[How we will stop doing this]"]])
    ppt_table_slide(prs, "Start — What We Should Try",
                    ["Item", "Votes", "Owner", "Try by"],
                    [["[New idea or experiment]", "[N]", "[Name]", "[Date]"],
                     ["[New idea or experiment]", "[N]", "[Name]", "[Date]"]])
    ppt_table_slide(prs, "Actions",
                    ["Action", "Owner", "Due Date"],
                    [["[Action from retro]", "[Name]", "[Date]"],
                     ["[Action from retro]", "[Name]", "[Date]"],
                     ["[Previous action — status check]", "[Name]", "Done / In Progress"]])
    prs.save(path)

# ── 10. DECOMMISSION ─────────────────────────────────────────────────────────

def decommission_plan(path):
    doc = Document()
    doc_heading(doc, "Decommission Plan")
    doc_field(doc, "Product / Service", "")
    doc_field(doc, "Date", "")
    doc_field(doc, "Owner", "")
    doc_field(doc, "Target Decommission Date", "")
    doc.add_paragraph()
    doc_heading(doc, "Reason for Decommission", 2)
    doc_checklist(doc, [
        "End of product life — superseded by [replacement]",
        "Low usage — does not justify ongoing operating cost",
        "Strategic change — no longer aligned with direction",
        "Technical — cannot be maintained sustainably",
        "Other: [Describe]",
    ])
    doc_heading(doc, "Impact Assessment", 2)
    doc_heading(doc, "Users and Customers", 3)
    doc_table(doc,
              ["Segment", "Number Affected", "Current Usage", "Migration Path"],
              [["[Segment]", "[N]", "[Description]", "[What they should use instead]"]])
    doc.add_paragraph()
    doc_heading(doc, "Dependent Systems", 3)
    doc_table(doc,
              ["System", "Owner", "Integration Type", "Action Required"],
              [["[System]", "[Team]", "[API / Data feed]", "[Migrate / Update / Remove]"]])
    doc.add_paragraph()
    doc_heading(doc, "Data", 3)
    doc_table(doc,
              ["Dataset", "Retention Requirement", "Action", "Owner"],
              [["[Dataset]", "[X years / None]", "Archive / Delete / Migrate", "[Name]"]])
    doc.add_paragraph()
    doc_heading(doc, "Communication Plan", 2)
    doc_table(doc,
              ["Audience", "Message", "Channel", "Date"],
              [["[Audience]", "[Key message]", "[Email / In-app / Blog]", "[Date]"]])
    doc.add_paragraph()
    doc_heading(doc, "Decommission Steps", 2)
    doc_table(doc,
              ["Step", "Action", "Owner", "Target Date", "Done"],
              [["1", "Final comms sent to affected users", "[Name]", "[Date]", "☐"],
               ["2", "New inflow blocked", "[Name]", "[Date]", "☐"],
               ["3", "Data archived or migrated", "[Name]", "[Date]", "☐"],
               ["4", "Service removed from production", "[Name]", "[Date]", "☐"],
               ["5", "DNS / routing cleaned up", "[Name]", "[Date]", "☐"],
               ["6", "Monitoring and alerting removed", "[Name]", "[Date]", "☐"],
               ["7", "Documentation marked as archived", "[Name]", "[Date]", "☐"],
               ["8", "Post-decommission review completed", "[Name]", "[Date]", "☐"]])
    doc.save(path)

def migration_checklist(path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Migration Checklist"
    xl_title(ws, "Migration Checklist",
             "Step-by-step checklist for migrating users and data to a replacement system")
    ws.row_dimensions[3].height = 6
    for r, (label, val) in enumerate([("Product Being Retired", "[Name]"),
                                       ("Migration Target", "[Replacement / None]"),
                                       ("Date", "[DD/MM/YYYY]"),
                                       ("Owner", "[Name, Role]")], start=4):
        ws.cell(r, 1).value = label; ws.cell(r, 1).font = xl_font(bold=True)
        ws.cell(r, 2).value = val;   ws.cell(r, 2).font = xl_font(color=GRAY)
    ws.row_dimensions[8].height = 10

    def checklist_block(start, title, items):
        xl_section(ws, start, title, span=4)
        xl_header(ws, start + 1,
                  [("Done", 8), ("Item", 60), ("Owner", 20), ("Notes", 28)])
        for i, item in enumerate(items):
            xl_row(ws, start + 2 + i, ["☐", item, "[Name]", ""], (start + 2 + i) % 2 == 0)
        return start + 2 + len(items) + 1

    row = checklist_block(9, "Pre-Migration — Planning", [
        "Migration scope defined",
        "Migration approach agreed",
        "Timeline communicated to all affected parties",
        "Rollback plan defined",
    ])
    row = checklist_block(row, "Pre-Migration — User Readiness", [
        "Users notified [X weeks] in advance",
        "Self-service migration guide published",
        "Support channel established",
        "High-value users contacted individually",
    ])
    row = checklist_block(row, "Migration Execution", [
        "Data export from source system",
        "Exported data completeness validated",
        "Data imported to target system",
        "Imported data validated in target",
        "Integrations updated to point to target",
        "User traffic redirected to target",
    ])
    row = checklist_block(row, "Post-Migration", [
        "Data integrity verified in target",
        "No critical data gaps identified",
        "Source system access revoked",
        "Users confirmed as successfully migrated",
        "Support tickets related to migration resolved",
        "Migration completion communicated",
    ])
    ws.sheet_view.showGridLines = False
    wb.save(path)


# ════════════════════════════════════════════════════════════════════════════
#  RUN ALL
# ════════════════════════════════════════════════════════════════════════════

BASE = os.path.dirname(os.path.abspath(__file__))

files = {
    "templates/inflow": [
        ("inflow-request-template.docx", inflow_request),
        ("stakeholder-brief.docx",        stakeholder_brief),
    ],
    "templates/value-proposition": [
        ("value-proposition-canvas.docx",  vp_canvas),
        ("success-metrics-template.xlsx",  success_metrics),
    ],
    "templates/product-discovery": [
        ("research-plan-template.docx", research_plan),
        ("assumption-log.xlsx",          assumption_log),
    ],
    "templates/align-and-plan": [
        ("delivery-plan-template.xlsx", delivery_plan),
        ("risk-register.xlsx",           risk_register),
    ],
    "templates/develop-and-deliver": [
        ("definition-of-done.docx",      definition_of_done),
        ("sprint-review-template.pptx",  sprint_review_pptx),
    ],
    "templates/deploy": [
        ("deployment-checklist.xlsx",   deployment_checklist),
        ("rollback-plan-template.docx", rollback_plan),
    ],
    "templates/operate-and-monitor": [
        ("incident-report-template.docx",    incident_report),
        ("operational-review-template.pptx", operational_review_pptx),
    ],
    "templates/measure-and-learn": [
        ("measurement-plan-template.xlsx", measurement_plan),
        ("insights-report-template.pptx", insights_report_pptx),
    ],
    "templates/adapt-and-optimize": [
        ("optimization-backlog-template.xlsx", optimization_backlog),
        ("retrospective-template.pptx",        retrospective_pptx),
    ],
    "templates/decommission": [
        ("decommission-plan-template.docx", decommission_plan),
        ("migration-checklist.xlsx",         migration_checklist),
    ],
}

for folder, items in files.items():
    full_folder = os.path.join(BASE, folder)
    ensure(full_folder)
    for filename, fn in items:
        path = os.path.join(full_folder, filename)
        fn(path)
        print(f"  ✓  {folder}/{filename}")

print("\nAll templates generated.")
